import pandas as pd
import json
import xml.etree.ElementTree as ET
import re
from pypdf import PdfReader
from pptx import Presentation
import io
from docx import Document

def parse_file(uploaded_file):
    """
    Parses the uploaded file based on its extension or content.
    Returns a DataFrame with columns: ['Severity', 'Name', 'Description', 'Category', 'File_Location']
    """
    filename = uploaded_file.name
    results = []

    # Helper to standardize columns
    def create_df(data):
        df = pd.DataFrame(data, columns=['Severity', 'Name', 'Description', 'Category', 'File_Location'])
        return df

    try:
        if filename.endswith('.xml') or filename.endswith('.nessus'):
            # Try Nessus first
            try:
                tree = ET.parse(uploaded_file)
                root = tree.getroot()
                
                # Check if it's Nessus (ReportHost/ReportItem structure)
                if root.tag == 'NessusClientData_v2':
                    for report_host in root.findall('.//ReportHost'):
                        for report_item in report_host.findall('ReportItem'):
                            severity_val = int(report_item.get('severity'))
                            severity_map = {4: 'Critical', 3: 'High', 2: 'Medium', 1: 'Low', 0: 'Info'}
                            
                            results.append({
                                'Severity': severity_map.get(severity_val, 'Info'),
                                'Name': report_item.get('pluginName'),
                                'Description': report_item.findtext('description', default=''),
                                'Category': report_item.get('pluginFamily', 'Nessus'),
                                'File_Location': report_host.get('name')
                            })
                    return create_df(results)
                
                # Check if specific Burp XML structure (issues)
                # Burp XML export usually has <issues><issue>...
                if root.tag == 'issues':
                    for issue in root.findall('issue'):
                        results.append({
                            'Severity': issue.findtext('severity', default='Info'),
                            'Name': issue.findtext('name', default='Unknown'),
                            'Description': issue.findtext('issueDetail', default='') or issue.findtext('issueBackground', default=''),
                            'Category': 'Burp Suite',
                            'File_Location': issue.findtext('host', default='') + issue.findtext('path', default='')
                        })
                    return create_df(results)

            except Exception:
                # Fallback or invalid XML
                pass

        elif filename.endswith('.json'):
            # OWASP ZAP or generic JSON
            content = json.load(uploaded_file)
            # OWASP ZAP Structure: {"site": [{"alerts": [...]}]}
            if 'site' in content:
                for site in content['site']:
                    for alert in site.get('alerts', []):
                        results.append({
                            'Severity': alert.get('riskdesc', 'Info').split(' ')[0], # "High (Medium Confidence)" -> High
                            'Name': alert.get('alert', 'Unknown'),
                            'Description': alert.get('desc', ''),
                            'Category': 'Vulnerability',
                            'File_Location': site.get('@name', '')
                        })
                return create_df(results)

        elif filename.endswith(('.py', '.js', '.txt', '.java', '.c', '.cpp')):
            # Basic SAST
            code_content = uploaded_file.getvalue().decode('utf-8', errors='ignore')
            lines = code_content.split('\n')
            
            patterns = {
                'eval(': 'High',
                'exec(': 'High',
                'password =': 'Medium',
                'password=': 'Medium', # tighter match
                'TODO': 'Info',
                'FIXME': 'Low',
                'hardcoded': 'Medium'
            }

            for i, line in enumerate(lines):
                for pattern, severity in patterns.items():
                    if pattern in line:
                         results.append({
                            'Severity': severity,
                            'Name': f'Potentially Unsafe Pattern: {pattern}',
                            'Description': f'Found pattern "{pattern}" at line {i+1}',
                            'Category': 'Static Analysis',
                            'File_Location': f'{filename}:{i+1}'
                        })
            return create_df(results)

        elif filename.endswith(('.pdf', '.pptx', '.ppt', '.docx', '.doc')):
            # Deep Document Scan (PDF, PPT, Word)
            text_content = ""
            if filename.endswith('.pdf'):
                try:
                    pdf_reader = PdfReader(uploaded_file)
                    for page in pdf_reader.pages:
                         text_content += page.extract_text() + "\n"
                except Exception as e:
                     return pd.DataFrame([{'Severity': 'Info', 'Name': 'PDF Parse Error', 'Description': str(e), 'Category': 'Error', 'File_Location': filename}])
            elif filename.endswith(('.pptx', '.ppt')):
                try:
                    ppt = Presentation(uploaded_file)
                    for slide in ppt.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_content += shape.text + "\n"
                except Exception as e:
                     return pd.DataFrame([{'Severity': 'Info', 'Name': 'PPT Parse Error', 'Description': str(e), 'Category': 'Error', 'File_Location': filename}])
            elif filename.endswith(('.docx', '.doc')):
                try:
                    if filename.endswith('.doc'):
                        # Warning: python-docx doesn't support .doc (binary). We'll try but likely fail or need simple text read.
                        # For now, treat as binary/text or fail gracefully.
                        pass 
                    doc = Document(uploaded_file)
                    for para in doc.paragraphs:
                        text_content += para.text + "\n"
                except Exception as e:
                     return pd.DataFrame([{'Severity': 'Info', 'Name': 'Word Parse Error', 'Description': f"Could not parse Word doc: {str(e)}", 'Category': 'Error', 'File_Location': filename}])

            # Deep Analysis
            # 1. Regex for CVEs
            cve_pattern = r'CVE-\d{4}-\d+'
            cves = re.findall(cve_pattern, text_content)
            for cve in set(cves): # Deduplicate
                 results.append({
                    'Severity': 'High',
                    'Name': f'Known Vulnerability: {cve}',
                    'Description': f'Detected Common Vulnerability Exposure ID {cve}',
                    'Category': 'CVE Detection',
                    'File_Location': filename
                })

            # 2. Context-Aware Keyword Search (Sentence-level)
            # Split text into rough sentences (by period or newline)
            sentences = re.split(r'[.\n]', text_content)
            keywords = {
                'SQL Injection': 'Critical',
                'XSS': 'High',
                'Remote Code Execution': 'Critical',
                'Hardcoded Password': 'High',
                'Confidential': 'Medium',
                'Private Key': 'Critical',
            }
            
            for sent in sentences:
                sent = sent.strip()
                if len(sent) > 200: continue # Skip massive blocks
                sent_lower = sent.lower()
                for key, severity in keywords.items():
                    if key.lower() in sent_lower:
                        # Found a keyword in a sentence
                         results.append({
                            'Severity': severity,
                            'Name': f'Potential Issue: {key}',
                            'Description': f'Context: "...{sent}..."', # Show the context
                            'Category': 'Deep Document Scan',
                            'File_Location': filename
                        })
            
            if not results:
                 results.append({
                            'Severity': 'Info',
                            'Name': 'Scan Complete',
                            'Description': 'No deep security indicators found in document text.',
                            'Category': 'Scan Info',
                            'File_Location': filename
                        })
            return create_df(results)

        elif filename.endswith(('.csv', '.xlsx')):
             # Generic Table
            if filename.endswith('.csv'):
                df = pd.read_csv(uploaded_file)
            else:
                df = pd.read_excel(uploaded_file)
            
            # Fuzzy match columns
            col_map = {}
            for col in df.columns:
                lower_col = col.lower()
                if 'severity' in lower_col or 'risk' in lower_col:
                    col_map['Severity'] = col
                elif 'name' in lower_col or 'vuln' in lower_col or 'title' in lower_col:
                    col_map['Name'] = col
                elif 'desc' in lower_col:
                    col_map['Description'] = col
                elif 'cat' in lower_col: # category
                     col_map['Category'] = col
                elif 'loc' in lower_col or 'path' in lower_col or 'url' in lower_col or 'file' in lower_col:
                    col_map['File_Location'] = col

            normalized_results = []
            for _, row in df.iterrows():
                normalized_results.append({
                    'Severity': row[col_map['Severity']] if 'Severity' in col_map else 'Info',
                    'Name': row[col_map['Name']] if 'Name' in col_map else 'Unknown',
                    'Description': row[col_map['Description']] if 'Description' in col_map else '',
                    'Category': row[col_map['Category']] if 'Category' in col_map else 'Generic',
                    'File_Location': row[col_map['File_Location']] if 'File_Location' in col_map else ''
                })
            return create_df(normalized_results)

    except Exception as e:
        # Return empty DF on error but log it?
        # For now just return empty with error as description
        return pd.DataFrame([{'Severity': 'Info', 'Name': 'Parse Error', 'Description': str(e), 'Category': 'Error', 'File_Location': ''}])

    return create_df([])
